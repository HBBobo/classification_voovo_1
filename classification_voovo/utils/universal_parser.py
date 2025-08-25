import os
from typing import List
import pdfplumber
from docx import Document
from pptx import Presentation

from utils.paragraph_splitter import split_text_into_paragraphs_with_gemini_async

async def _get_paragraphs_from_pdf_async(file_path: str) -> List[str]:
    """Extracts raw text from a PDF and uses Gemini to split it into paragraphs, asynchronously."""
    print(f"  -> Using PDF method with Gemini splitter...")
    try:
        with pdfplumber.open(file_path) as pdf:
            raw_text = " ".join(
                page.extract_text(x_tolerance=1).replace('\n', ' ')
                for page in pdf.pages if page.extract_text()
            )
        return await split_text_into_paragraphs_with_gemini_async(raw_text)
    except Exception as e:
        print(f"    -> Error processing PDF {file_path}: {e}")
        return []

def _get_paragraphs_from_docx(file_path: str) -> List[str]:
    """Extracts paragraphs directly from a .docx file's structure and filters out empty ones."""
    print("  -> Using DOCX method (direct structure read)...")
    try:
        doc = Document(file_path)
        return [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    except Exception as e:
        print(f"    -> Error processing DOCX {file_path}: {e}")
        return []

def _get_paragraphs_from_pptx(file_path: str) -> List[str]:
    """Extracts text from each slide of a .pptx file, filters out empty ones."""
    print("  -> Using PPTX method (one text chunk per slide)...")
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            
            slide_text = " ".join(text_runs).strip()
            if slide_text:
                slide_texts.append(slide_text)
        return slide_texts
    except Exception as e:
        print(f"    -> Error processing PPTX {file_path}: {e}")
        return []

async def extract_paragraphs_from_file_async(file_path: str) -> List[str]:
    """
    Universally extracts paragraphs from PDF, DOCX, or PPTX files.

    This function acts as a dispatcher, checking the file extension and
    calling the appropriate specialized function to handle the file type.

    Args:
        file_path: The full path to the file.

    Returns:
        A list of strings, where each string is a paragraph or a
        logical chunk of text from the document.
    """
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return []

    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    
    print(f"\nProcessing file: '{os.path.basename(file_path)}'")

    if extension == '.pdf':
        return await _get_paragraphs_from_pdf_async(file_path)
    elif extension == '.docx':
        return _get_paragraphs_from_docx(file_path)
    elif extension == '.pptx':
        return _get_paragraphs_from_pptx(file_path)
    else:
        print(f"  -> Unsupported file type: '{extension}'. Skipping.")
        return []

if __name__ == "__main__":
    # This block is for direct testing of the module
    import asyncio
    async def test_parser():
        files_to_process = [
            "data/content/Jagellók.docx" 
        ]

        all_extracted_content = {}
        for file in files_to_process:
            paragraphs = await extract_paragraphs_from_file_async(file)
            if paragraphs:
                all_extracted_content[file] = paragraphs
                print(f"  -> Successfully extracted {len(paragraphs)} paragraph(s).")
                print(f"  -> Sample: '{paragraphs[0]}'")

        print("\n--- Universal Parser Demonstration Complete ---")

    asyncio.run(test_parser())